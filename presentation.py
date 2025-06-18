from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import os
import platform



def create_ppt(filename="output.pptx"):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.add_picture("tMobile.jpg", Inches(0), Inches(0), width=Inches(2), height=Inches(2))
    slide.shapes.title.text = "WELCOME TO T-MOBILE"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb=RGBColor(226, 0, 116)
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.name='Times New Roman'
    content=("T-Mobile US, Inc. (NASDAQ: TMUS) is America’s supercharged Un-carrier, delivering an advanced "
        "4G LTE and transformative nationwide 5G network that will offer reliable connectivity for all. "
        "T-Mobile’s customers benefit from its unmatched combination of value and quality, unwavering "
        "obsession with offering them the best possible service experience and undisputable drive for disruption "
        "that creates competition and innovation in wireless and beyond.")
    placeholder=slide.placeholders[1]
    placeholder.text=content

    para=placeholder.text_frame.paragraphs[0]
    run=para.runs[0]
    run.font.size=Pt(11)
    run.font.name='Times New Roman'

    #stock data slide



    df=pd.read_excel("tmobileStock.xlsx")
    data_slide= prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(2.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)

    title_box = data_slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "STOCK ANALYSIS"
    run.font.color.rgb = RGBColor(226, 0, 116)
    run.font.name = 'Times New Roman'
    run.font.size=Pt(44)

    rows, cols= df.shape[0]+1, df.shape[1]
    left=Inches(0.5)
    top=Inches(2)
    width=Inches(9)
    height=Inches(1.5)
    table=data_slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i, col_name in enumerate(df.columns):
        cell=table.cell(0, i)
        cell.text=str(col_name)

        paragraph= cell.text_frame.paragraphs[0]
        run=paragraph.runs[0]
        run.font.bold=True
        run.font.name='Times New Roman'
        run.font.size=Pt(20)
        run.font.color.rgb=RGBColor(255, 255, 255)

        cell.fill.solid()
        cell.fill.fore_color.rgb=RGBColor(226, 0, 116)

    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            cell=table.cell(i+1, j)
            cell.text=str(df.iat[i,j])

            paragraph = cell.text_frame.paragraphs[0]
            run=paragraph.runs[0]
            run.font.size=Pt(11)
            run.font.name='Times New Roman'

            if i%2==0:
                cell.fill.solid()
                cell.fill.fore_color.rgb=RGBColor(245, 171, 209)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb=RGBColor(255, 255, 255)


    #Thankyou slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background=slide.background
    fill=background.fill
    fill.solid()
    fill.fore_color.rgb=RGBColor(226, 0, 116)
    slide.shapes.title.text = "THANKYOU"
    title=slide.shapes.title
    title.width=Inches(6)
    title.height=Inches(1)
    title.left=Inches(2)
    title.top=Inches(3)


    paragraph=slide.shapes.title.text_frame.paragraphs[0]
    run=paragraph.runs[0]
    run.font.color.rgb=RGBColor(255, 255, 255)
    paragraph.alignment=PP_ALIGN.CENTER
    run.font.name='Times New Roman'
    

    prs.save(filename)
    print(f"Saved presentation to {filename}")
    open_ppt(filename)

def open_ppt(filename):
    if platform.system() == "Darwin":       # macOS
        os.system(f"open {filename}")
    elif platform.system() == "Windows":    # Windows
        os.system(f"start {filename}")
    else:                                   # Linux
        os.system(f"xdg-open {filename}")

if __name__ == "__main__":
    create_ppt()
